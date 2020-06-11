import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt
import pptx_sp as ps
from pathlib import Path
import os

print("------------------------------------------------------------------ NEW COMMAND ------------------------------------------------------------------")
pa = os.getcwd()
spinput = os.path.join(pa, "database/input/", "sp-ptc1-2020.xlsx")
spfile = pd.ExcelFile(spinput)
sptitle = Path(spinput).stem
sheetname = spfile.sheet_names
sheetcount = len(sheetname)
print("Number of sheets: %d" % sheetcount)
df_list = []
df_idlist = []

for s in range(sheetcount):
    df_ = pd.read_excel(spinput, sheet_name=s, skiprows=2,
                        usecols='F:K', nrows=None)
    df_ = df_.iloc[:, [1, 3, 5]]
    df_list.append(df_)
    df_id = pd.read_excel(spinput, sheet_name=s, skiprows=2,
                          usecols='A:C', nrows=None)
    df_idlist.append(df_id)

output1 = os.path.join(pa, "database/output/", 'sp_output1.xlsx')
ps.write_excel(df_list, output1, sheetname)
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_content_layout = prs.slide_layouts[1]
sl_title, sh_title = ps.add_slide(prs, title_slide_layout, sptitle)

for i, df in enumerate(df_list):
    # add table here and image here
    df = ps.reindex_column(df)
    print("-----------------------------------------------------\n{}\nSheet/df shape: {}, rowcount: {}".format(
        sheetname[i], str(df.shape), df.shape[0]))
    df, sheet = ps.normalize_df(df, 0)
    slide, shape = ps.add_slide(prs, title_content_layout, sheetname[i])
    print(sheet)
    ps.create_table(slide, sheet, 1, 10)
    # pic = shape.add_picture('image.jpg', Inches(1), Inches(1))
    df_list[i] = df  # replace df with normalized df
    dd = df_idlist[i]
    dd = ps.reindex_column(dd)
    print(dd)
    dd, sh = ps.normalize_df(dd, 1)
    print(sh)
    ps.insert_image(pa, shape, sh)

# output2 = os.path.join(pa, "database/output/", 'sp_output2.xlsx')
# output3 = os.path.join(pa, "database/output/", 'sp_output3.xlsx')
# createsp = os.path.join(pa, "database/output/", "createsp.pptx")
# ps.write_excel(df_list, output2, sheetname)
# ps.write_excel(df_idlist, output3, sheetname)
# sl_end, sh_end = ps.add_slide(prs, title_slide_layout, "End")
# prs.save(createsp)